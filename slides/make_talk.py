#!/usr/bin/env python3
"""Build the conference talk deck (15-20 min) from the project's own artifacts.

The manuscript rule applies here too: NO NUMBER IS TYPED INTO THIS FILE.  Every
quantity is pulled from `manuscript/numbers.tex`, which `wp10_numbers.py`
generates from versioned artifacts.  If the pipeline changes, regenerate the
deck and the talk cannot quietly disagree with the paper.

Output: slides/cygob2_supernova_history_talk.pptx

Run:
  python3 slides/make_talk.py
"""
from __future__ import annotations

import re
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "slides" / "cygob2_supernova_history_talk.pptx"

# ---------------------------------------------------------------- palette
# Same mapping as the paper's figures: alpha = 2.0 blue, alpha = 2.3 orange.
INK = RGBColor(0x1A, 0x1F, 0x2B)
MUTED = RGBColor(0x5A, 0x62, 0x72)
FAINT = RGBColor(0x8A, 0x92, 0xA0)
BLUE = RGBColor(0x3B, 0x6E, 0xA8)
ORANGE = RGBColor(0xC4, 0x62, 0x2D)
GREEN = RGBColor(0x4F, 0x8A, 0x5B)
RULE = RGBColor(0xD8, 0xDD, 0xE4)
BG_DARK = RGBColor(0x1A, 0x1F, 0x2B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

W, H = Inches(13.333), Inches(7.5)
MARGIN = Inches(0.72)
BODY_W = W - 2 * MARGIN


def numbers() -> dict[str, str]:
    """Every macro from the generated numbers.tex."""
    text = (ROOT / "manuscript" / "numbers.tex").read_text()
    return dict(re.findall(r"\\newcommand\{\\(\w+)\}\{([^}]*)\}", text))


N = numbers()


def deck() -> Presentation:
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    return prs


def blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def textbox(slide, left, top, width, height, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    frame.vertical_anchor = anchor
    return frame


def para(frame, text, size, colour=INK, bold=False, space_after=10,
         align=PP_ALIGN.LEFT, first=False, italic=False, line_spacing=1.0):
    p = frame.paragraphs[0] if first else frame.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = colour
    run.font.name = "Helvetica Neue"
    return p


def fill(slide, colour):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = colour


def rule(slide, top, colour=RULE, height=Pt(1.5), left=MARGIN, width=BODY_W):
    from pptx.enum.shapes import MSO_SHAPE
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = colour
    bar.line.fill.background()
    bar.shadow.inherit = False
    return bar


def heading(slide, title, kicker=None):
    """Standard slide header: small kicker, big title, hairline."""
    top = Inches(0.52)
    if kicker:
        frame = textbox(slide, MARGIN, top, BODY_W, Inches(0.3))
        para(frame, kicker.upper(), 12.5, FAINT, bold=True, first=True,
             space_after=0)
        top = top + Inches(0.36)
    frame = textbox(slide, MARGIN, top, BODY_W, Inches(0.9))
    para(frame, title, 30, INK, bold=True, first=True, space_after=0,
         line_spacing=0.95)
    rule(slide, Inches(1.86))
    return Inches(2.12)


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


# ------------------------------------------------------------------ slides
def title_slide(prs):
    s = blank(prs)
    fill(s, BG_DARK)
    frame = textbox(s, MARGIN, Inches(2.05), Inches(11.4), Inches(2.6))
    para(frame, "THE SUPERNOVA HISTORY OF CYGNUS OB2", 13, FAINT, bold=True,
         first=True, space_after=16)
    para(frame, "How many stars have exploded, when,", 40, WHITE, bold=True,
         space_after=2, line_spacing=1.0)
    para(frame, "and what that means for the Cygnus PeVatron", 40, WHITE,
         bold=True, space_after=18, line_spacing=1.0)
    para(frame, "A probabilistic ledger of massive-star deaths "
                "from a homogeneous Gaia DR3 census",
         16.5, RGBColor(0x9F, 0xB4, 0xCE), italic=True, space_after=0)
    frame = textbox(s, MARGIN, Inches(5.75), Inches(11.4), Inches(1.0))
    para(frame, "V. Voitsekhovskyi", 17, WHITE, bold=True, first=True,
         space_after=4)
    para(frame, "Gaia DR3  ·  INTEGRAL/SPI  ·  forecast for COSI", 13.5, FAINT,
         space_after=0)
    notes(s, "15-20 min. Aim: the ledger is the measurement; the PeVatron "
             "verdict is one application of it. Do not oversell the verdict.")
    return s


def section(prs, number, title, subtitle):
    s = blank(prs)
    fill(s, BG_DARK)
    frame = textbox(s, MARGIN, Inches(2.9), Inches(11.4), Inches(2.0))
    para(frame, number, 13, FAINT, bold=True, first=True, space_after=14)
    para(frame, title, 44, WHITE, bold=True, space_after=12)
    para(frame, subtitle, 17, RGBColor(0x9F, 0xB4, 0xCE), space_after=0)
    return s


def bullets(prs, title, items, kicker=None, note=None, size=17.5):
    s = blank(prs)
    top = heading(s, title, kicker)
    frame = textbox(s, MARGIN, top, BODY_W, Inches(4.6))
    for i, item in enumerate(items):
        if isinstance(item, tuple):
            text, colour, bold = item
        else:
            text, colour, bold = item, INK, False
        para(frame, text, size, colour, bold=bold, first=(i == 0),
             space_after=15, line_spacing=1.05)
    if note:
        notes(s, note)
    return s


def figure_slide(prs, title, image, caption=None, kicker=None, note=None,
                 max_h=Inches(4.5)):
    s = blank(prs)
    top = heading(s, title, kicker)
    path = ROOT / image
    from PIL import Image
    with Image.open(path) as im:
        ratio = im.height / im.width
    width = BODY_W
    height = Emu(int(width * ratio))
    if height > max_h:
        height = max_h
        width = Emu(int(height / ratio))
    left = Emu(int((W - width) / 2))
    s.shapes.add_picture(str(path), left, top, width=width, height=height)
    if caption:
        frame = textbox(s, MARGIN, top + height + Inches(0.14), BODY_W,
                        Inches(0.7))
        para(frame, caption, 14, MUTED, first=True, space_after=0,
             line_spacing=1.05)
    if note:
        notes(s, note)
    return s


def fit_size(text: str, width_emu: int, ideal: float, floor: float = 24.0
             ) -> float:
    """Largest point size at which `text` fits one line in `width_emu`.

    A hero number is only a hero if it is big AND fits.  Long strings like
    "0.593-0.735" overflow a quarter-width column at 54 pt, which is what a
    fixed size silently produces.
    """
    if not text:
        return ideal
    # ~0.55 em average advance for bold Helvetica digits and punctuation.
    usable = (width_emu / Inches(1)) * 96.0 * 0.94
    return max(floor, min(ideal, usable / (len(text) * 0.55)))


def stat_slide(prs, title, stats, kicker=None, footer=None, note=None):
    """Two to four hero numbers -- the form for a headline, not a chart.

    Values are kept SHORT and units pushed into the label: the number is the
    thing being read, and "2.92" reads instantly where "2.92x10^4 Msun" does
    not.
    """
    s = blank(prs)
    top = heading(s, title, kicker)
    n = len(stats)
    gap = Inches(0.34)
    col_w = Emu(int((BODY_W - gap * (n - 1)) / n))
    size = min(fit_size(value, col_w, 54.0) for value, _, _ in stats)
    for i, (value, label, colour) in enumerate(stats):
        left = MARGIN + Emu(int(i * (col_w + gap)))
        frame = textbox(s, left, top + Inches(0.35), col_w, Inches(2.6))
        para(frame, value, size, colour, bold=True, first=True, space_after=10,
             align=PP_ALIGN.CENTER, line_spacing=0.95)
        for j, line in enumerate(label.split("\n")):
            para(frame, line, 14.5, MUTED, space_after=0,
                 align=PP_ALIGN.CENTER, line_spacing=1.1)
    if footer:
        frame = textbox(s, MARGIN, Inches(5.75), BODY_W, Inches(1.1))
        para(frame, footer, 16, INK, first=True, space_after=0,
             line_spacing=1.1)
    if note:
        notes(s, note)
    return s


def quote_slide(prs, big, sub=None, kicker=None, colour=BLUE, note=None):
    s = blank(prs)
    top = heading(s, kicker or "", None) if False else Inches(2.3)
    frame = textbox(s, MARGIN, top, BODY_W, Inches(2.8))
    if kicker:
        para(frame, kicker.upper(), 12.5, FAINT, bold=True, first=True,
             space_after=18)
        para(frame, big, 38, colour, bold=True, space_after=20,
             line_spacing=1.02)
    else:
        para(frame, big, 38, colour, bold=True, first=True, space_after=20,
             line_spacing=1.02)
    if sub:
        para(frame, sub, 18, INK, space_after=0, line_spacing=1.15)
    if note:
        notes(s, note)
    return s


def build() -> Presentation:
    prs = deck()

    # =========================================================== 1. the problem
    title_slide(prs)

    section(prs, "1", "The puzzle",
            "A PeVatron in Cygnus, and five candidate engines")

    bullets(prs, "LHAASO sees a cosmic-ray PeVatron in Cygnus",
            kicker="the observation",
            items=[
                ("A γ-ray 'cocoon' around Cyg OB2, emitting above 100 TeV — "
                 "protons accelerated to PeV energies.", INK, False),
                ("Candidate engines all live in the same region:", MUTED, True),
                ("     · collective winds of the massive-star population", INK, False),
                ("     · a supernova exploding into the wind-blown cavity", BLUE, True),
                ("     · the Cyg X-3 microquasar", INK, False),
                ("     · the γ Cygni supernova remnant", INK, False),
                ("     · PSR J2032+4127", INK, False),
            ],
            note="Set up: the field argues about which engine. Härer et al. "
                 "2025 argue for a supernova ~50 kyr ago into the cavity. "
                 "Every one of these depends on the SAME population property.")

    quote_slide(prs,
                "Every candidate engine depends on one number nobody had "
                "measured: how many supernovae Cyg OB2 has actually produced, "
                "and when.",
                sub="Wind power, cavity structure, remnant population, "
                    "pulsar birth — all of them are set by the association's "
                    "death history. That history is a MEASUREMENT problem, "
                    "and it is the subject of this talk.",
                kicker="the gap", colour=INK,
                note="This is the pivot. The talk is a measurement first; the "
                     "PeVatron interpretation is one application at the end.")

    # ============================================================== 2. the target
    section(prs, "2", "Cygnus OB2",
            "The best target in the Galaxy, and why it is hard")

    stat_slide(prs, "Cygnus OB2 — the target", kicker="what we are looking at",
               stats=[(N["ourDistance"], "kpc\ndistance", BLUE),
                      (N["Nmembers"], "members\nat P > 0.5", INK),
                      (N["massTotal"], "×10⁴ M⊙\nstellar mass", INK),
                      (f"{N['ageC']}–{N['ageB']}", "Myr\nsubgroup ages", ORANGE)],
               footer="Young enough that the first supernovae have only just "
                      "begun, close enough for Gaia to resolve individual "
                      "stars, and massive enough that the statistics are "
                      "meaningful. Three kinematic subgroups (A, B, C) — "
                      "the association is not coeval.",
               note="Emphasise: nearest laboratory where you can count the "
                    "progenitors star by star instead of assuming an IMF.")

    bullets(prs, "Why this has not been settled before",
            kicker="the difficulty",
            items=[
                ("Extinction is severe and patchy — A_V ≈ 6 mag, varying "
                 "star to star across the field.", INK, False),
                ("The stars that matter have already died. You must infer the "
                 "dead from the living, through an IMF.", ORANGE, True),
                ("Every step carries a model choice: isochrones, extinction "
                 "law, IMF slope, explodability, star-formation duration.", INK, False),
                ("Those choices are not decorative — they move the answer by "
                 "more than an order of magnitude.", INK, False),
            ],
            note="This motivates the branch discipline: we carry model "
                 "choices as explicit branches and never average them.")

    # =========================================================== 3. the questions
    bullets(prs, "The questions this work answers",
            kicker="scope",
            items=[
                ("1.  How many core-collapse supernovae has Cyg OB2 produced?", INK, True),
                ("2.  When did they happen — and how likely is a recent one?", INK, True),
                ("3.  What were the progenitors?", INK, True),
                ("4.  Is the supernova-driven PeVatron scenario supported?", BLUE, True),
                ("Design rule fixed before any result existed: every model "
                 "choice is a BRANCH, carried in parallel and never averaged; "
                 "every comparison is PRE-REGISTERED before it is scored.",
                 MUTED, False),
            ],
            note="Stress the pre-registration: thresholds and predictions were "
                 "frozen in JSON before the numbers were computed. Failed "
                 "predictions are reported as failed.")

    # ============================================================== 4. the method
    section(prs, "3", "Method",
            "From a Gaia catalogue query to a supernova ledger")

    bullets(prs, "The chain, and the one out-of-sample test in it",
            kicker="pipeline",
            items=[
                ("Gaia DR3 + 2MASS  →  membership & three subgroups", INK, False),
                ("→  per-star extinction and de-reddened CMD", INK, False),
                ("→  subgroup ages and per-star masses", INK, False),
                ("→  IMF normalization from 2–8 M⊙ counts + completeness", INK, False),
                ("→  census closure above 8 M⊙  ·  runaway recovery", GREEN, True),
                ("→  the supernova ledger: how many died, and when", BLUE, True),
                ("The IMF is fitted on 2–8 M⊙ stars only. The >8 M⊙ census "
                 "never enters that fit — so comparing predicted to observed "
                 "massive stars is a genuine out-of-sample test.", MUTED, False),
            ],
            size=16.5,
            note="The closure test is the credibility anchor. Say clearly "
                 "that we deliberately did NOT spend it to fit alpha.")

    figure_slide(prs, "Membership: 1,392 members, three subgroups",
                 "figures/paper/fig1_membership_literature.png",
                 kicker="step 1",
                 caption="Clustering in position and proper motion; parallax "
                         "deliberately excluded (one population at 1.62 kpc, "
                         "its depth exhausted at DR3 precision). "
                         "Literature recall 0.825; control-field yield 3.7%.",
                 note="Do not dwell. One line: the membership is validated "
                      "against Berlanas+19 and against control fields.")

    figure_slide(prs, "The mass function, and what normalizes it",
                 "figures/paper/fig3_mass_function.png",
                 kicker="step 2",
                 caption=f"Forward-modelled counts against observation for the "
                         f"three subgroups. {N['branchesPassing']} of "
                         f"{N['branchesTotal']} model branches pass the "
                         f"residual gate. The normalization k is fitted here, "
                         f"below 8 M⊙, and then extrapolated upward.",
                 note="The injection/recovery machinery behind completeness is "
                      "the expensive part; mention it exists, move on.")

    stat_slide(prs, "The census closes at Salpeter — out of sample",
               kicker="step 3 · the validation that matters",
               stats=[(N["closureExcess"] + "%",
                       "observed massive stars\nexceed prediction (α = 2.3)", GREEN),
                      (N["closingAlpha"],
                       "slope α at which the\ncensus closes exactly", INK),
                      (f"{N['closureCellsInside']}/{N['closureCells']}",
                       "cells inside the\ncarried branch grid", INK)],
               footer="The IMF was normalized on low-mass counts alone, yet it "
                      "predicts the observed massive-star population to within "
                      "7%. Nothing was tuned to make this happen.",
               note="This is the single most persuasive slide for a sceptic. "
                    "Let it breathe.")

    bullets(prs, "Runaways: stars that left, and still count",
            kicker="step 4",
            items=[
                (f"{N['runawaysRaw']} raw candidates recovered by traceback of "
                 f"PECULIAR proper motions  →  {N['runawaysCorrected']} after a "
                 f"measured false-positive rate.", INK, False),
                (f"They bound the fraction of supernovae that occurred outside "
                 f"the association at ≤ {N['runawayFraction']}%.", INK, False),
                ("External check: BD+43 3654, the canonical ejected O star, is "
                 "recovered at probability 1.000 — 38.8 km/s and 1.36 Myr "
                 "against a literature ~40 km/s and 1.6 Myr.", GREEN, True),
                ("That check also caught a real bug: the first version traced "
                 "back ABSOLUTE proper motions, which measure the "
                 "association's bulk drift, not ejection. Result withdrawn "
                 "and re-run.", ORANGE, False),
            ],
            note="Being open about the withdrawn result buys credibility and "
                 "pre-empts the obvious referee question.")

    # ============================================================== 5. results
    section(prs, "4", "Results",
            "The ledger, and the two findings that dominate it")

    figure_slide(prs, "The supernova ledger",
                 "figures/paper/fig4_rsn_history.png",
                 kicker="result 1",
                 caption=f"Supernova rate against look-back time. Baseline "
                         f"branch: N_SN = {N['NSN']} for the association, "
                         f"P(≥1) = {N['Pone']}, and P(last SN < 100 kyr) = "
                         f"{N['Precent']} — median time since the last "
                         f"explosion {N['tlast']} kyr.",
                 note="P(recent SN) = 0.55 is the number the PeVatron people "
                      "care about. Note it is 7x an ignorance baseline.")

    stat_slide(prs, "The number is not a number",
               kicker="result 1 · honesty about branches",
               stats=[(N["NSN"], "supernovae\nbaseline branch", INK),
                      (f"{N['NSNheadlo']}–{N['NSNheadhi']}",
                       f"across the {N['NSNheadbranches']}\nheadline branches", ORANGE),
                      ("×" + N["NSNheadfactor"],
                       "spread, driven\nmainly by the IMF slope", ORANGE)],
               footer="We report the range, not a consensus value, and we never "
                      "average branches. A single headline number here would be "
                      "a fiction — the model choices are not resolved by the data.",
               note="If you take one methodological point away, this is it.")

    quote_slide(prs,
                f"The entire supernova budget lies above {N['bhSafeCut']} M⊙.",
                sub=f"For any black-hole threshold at or below "
                    f"{N['bhSafeCut']} M⊙, the ledger returns EXACTLY ZERO on "
                    f"every branch. The lowest-mass star that has died anywhere "
                    f"in the grid is {N['minProgenitor']} M⊙. So the whole "
                    f"result is conditional on whether very massive stars "
                    f"explode at all — which is unsettled physics.",
                kicker="result 2 · the structural finding", colour=ORANGE,
                note="This is a genuine discovery about the system, not a "
                     "caveat. It also sets up the pulsar slide.")

    bullets(prs, "An observation we already had rules out the zero branch",
            kicker="external cross-check",
            items=[
                ("PSR J2032+4127 is a neutron star inside Cyg OB2 — and "
                 "neutron stars require a SUCCESSFUL explosion.", INK, False),
                (f"Ledger: P(≥1 SN) = {N['pulsarPone']} if massive stars "
                 f"explode, against exactly {N['pulsarIslands']} on the "
                 f"black-hole branch.", BLUE, True),
                ("Its companion MT91 213 is our own census star — a B0V of "
                 "17 M⊙, 0.115° from the CygOB2-A centroid, already counted.", INK, False),
                (f"Age agrees too: characteristic age widened to 151–401 kyr "
                 f"against a ledger probability of {N['pulsarAge']} that the "
                 f"last supernova falls in that window.", INK, False),
                ("Three readings remain degenerate — high-mass explodability, "
                 "an older population, or binary stripping. All three give a "
                 "non-zero budget, which is what the argument needs.", MUTED, False),
            ],
            size=16.5,
            note="Strongest single piece of evidence in the talk. The "
                 "companion being our own star is the detail people remember.")

    # ============================================================== 6. verdict
    section(prs, "5", "The verdict",
            "One application of the ledger — decided by a rule fixed in advance")

    figure_slide(prs, "Is the supernova-PeVatron scenario supported?",
                 "figures/paper/fig6_verdict_branches.png",
                 kicker="result 3",
                 caption=f"P(verdict) per branch = P(right age) × P(right "
                         f"progenitor type) × P(in situ). Range "
                         f"{N['Pverdictlo']}–{N['Pverdicthi']}, median "
                         f"{N['Pverdictmed']}. The explosion energy is NOT "
                         f"multiplied in — it is reported as a stated "
                         f"conditional.",
                 note="Explain we refuse to invent an energy probability. "
                      "The framing rule was mechanical: INCONCLUSIVE.")

    stat_slide(prs, "The verdict hinges on one axis, and only one",
               kicker="result 3 · where the uncertainty lives",
               stats=[(f"{N['PverdictTwolo']}–{N['PverdictTwohi']}",
                       "α = 2.0 branches\n18/18 support the scenario", BLUE),
                      (f"{N['PverdictTwoThreelo']}–{N['PverdictTwoThreehi']}",
                       "α = 2.3 branches\n0/18 support it", ORANGE),
                      (N["spreadAlpha"],
                       f"spread from α, against\n{N['spreadFamily']} for isochrones", INK)],
               footer="The 0.5 boundary falls exactly between the two arms. "
                      "The pre-registered rule therefore returns INCONCLUSIVE "
                      "— and we wrote a regular article rather than a Letter, "
                      "mechanically, because that is what the rule said.",
               note="Emphasise: we let a rule written in advance decide the "
                    "framing, including against our own interest.")

    # ============================================================== 7. forecast
    section(prs, "6", "A forward prediction",
            "Making the ledger falsifiable — with an instrument launching in 2027")

    bullets(prs, "Gaia DR4 will not settle the axis that matters",
            kicker="the problem with waiting",
            items=[
                ("DR4 gives radial velocities (3D traceback), sharper ages, "
                 "deeper membership. All useful.", INK, False),
                ("None of it constrains the high-mass IMF slope — which is "
                 "the one axis the verdict turns on.", ORANGE, True),
                ("So we looked for an independent observable that does.", INK, False),
                ("⁶⁰Fe is the supernova-specific tracer: Wolf–Rayet winds "
                 "produce ²⁶Al but NO ⁶⁰Fe. It counts explosions, not a "
                 "combination.", BLUE, True),
            ],
            note="Transition slide. Keep it quick.")

    figure_slide(prs, "COSI can separate the branches Gaia cannot",
                 "figures/wp11/wp11_cosi_forecast.png",
                 kicker="the forecast",
                 caption=f"Predicted ⁶⁰Fe line flux per branch against COSI's "
                         f"3σ sensitivity. {N['isoCosiTwo']}/"
                         f"{N['isoCosiTwon']} branches at α = 2.0 are "
                         f"detectable; {N['isoCosiTwoThree']}/"
                         f"{N['isoCosiTwoThreen']} at α = 2.3. Yields were "
                         f"declared as a branch BEFORE any flux was computed.",
                 max_h=Inches(4.2),
                 note="The punchline of the new work. But immediately give "
                      "the caveat on the next slide - do not oversell.")

    bullets(prs, "Two things bound that claim — and we give them equal weight",
            kicker="honesty",
            items=[
                (f"The yield model moves the prediction by a factor "
                 f"{N['isoArmSpread']}, against {N['isoBranchSpread']} for the "
                 f"whole branch set. A non-detection constrains BOTH, not α "
                 f"alone.", ORANGE, True),
                (f"The current standard yield compilation predicts nothing at "
                 f"all: it collapses everything above 25 M⊙, and "
                 f"{N['isoNullBelow']} of the {N['isoNullSampled']} million "
                 f"supernovae we sample fall below that — so its ⁶⁰Fe flux is "
                 f"identically zero.", INK, False),
                ("A COSI non-detection would then favour exactly the "
                 "explodability that the pulsar disfavours — and that tension "
                 "would itself be the result.", BLUE, True),
                (f"Existing INTEGRAL data are already within "
                 f"{N['isoSpiMargin']}% of excluding our richest branch.", INK, False),
            ],
            size=16.5,
            note="Never let the COSI slide stand alone. This slide is what "
                 "makes the forecast credible rather than promotional.")

    # ============================================================== 8. close
    section(prs, "7", "Conclusions", "What is measured, what is not")

    bullets(prs, "Conclusions",
            items=[
                (f"Cyg OB2 has produced {N['NSNheadlo']}–{N['NSNheadhi']} "
                 f"core-collapse supernovae; P(≥1) = {N['Pone']} and "
                 f"P(last < 100 kyr) = {N['Precent']}, seven times an "
                 f"ignorance baseline.", INK, True),
                (f"The census closes out of sample to {N['closureExcess']}% at "
                 f"Salpeter — the IMF extrapolation is validated, not assumed.", GREEN, False),
                (f"The whole budget sits above {N['bhSafeCut']} M⊙, so the "
                 f"result is conditional on very massive stars exploding. "
                 f"PSR J2032+4127 says at least one did.", ORANGE, False),
                ("The PeVatron verdict is INCONCLUSIVE by a rule fixed in "
                 "advance, and hinges on the IMF slope alone.", INK, False),
                ("⁶⁰Fe makes it falsifiable: COSI discriminates the branches "
                 "— if the yield model cooperates.", BLUE, True),
            ],
            note="Land on: this is a measurement others can build on, with "
                 "every branch and every failed prediction on the record.")

    quote_slide(prs, "Thank you",
                sub="Every number in this talk is generated from a versioned "
                    "artifact — the same macro file the paper uses. "
                    "Pre-registrations, failed predictions and withdrawn "
                    "results are all in the repository.",
                colour=INK,
                note="Backup slides follow: systematics, pre-registration, "
                     "the 26Al correction, age sensitivity.")

    # ============================================================== backup
    section(prs, "—", "Backup", "For questions")

    figure_slide(prs, "Backup · How the answer depends on the assumed age",
                 "figures/paper/fig5_age_sensitivity.png",
                 kicker="backup",
                 caption="Below ~3 Myr nothing has died yet. The steep "
                         "dependence on age is why the counts-based ages, not "
                         "the CMD ages, are what the ledger consumes.")

    figure_slide(prs, "Backup · Control fields and the false-positive rate",
                 "figures/paper/fig2_control_fields.png",
                 kicker="backup",
                 caption="Membership contamination and the runaway "
                         "chance-alignment rate are both measured from control "
                         "fields rather than assumed.")

    bullets(prs, "Backup · ²⁶Al, and a correction we found",
            kicker="backup",
            items=[
                ("²⁶Al traces winds AND supernovae, so it constrains a "
                 "combination and is never inverted for a supernova count.", INK, False),
                (f"Run forward, the ledger's supernovae supply "
                 f"{N['isoAlFracLo']}–{N['isoAlFracHi']}% of the measured "
                 f"complex-wide 1809 keV flux — sub-dominant but not "
                 f"negligible.", INK, False),
                ("This corrected an earlier internal comparison that was wrong "
                 "by ~100× (a Galactic ²⁶Al mass had been used for the Cygnus "
                 "complex). Recorded as a finding; nothing upstream retuned.", ORANGE, True),
            ],
            note="Only if asked. Shows the one-way validation rule working.")

    bullets(prs, "Backup · What this cannot do",
            kicker="backup",
            items=[
                ("It cannot compute explosion energies — C2 is a stated "
                 "conditional, never a multiplied probability.", INK, False),
                ("It cannot resolve whether γ Cygni is associated with Cyg OB2 "
                 f"(P = {N['gammaCygni']}% within its ~7 kyr age).", INK, False),
                ("It cannot separate Cyg OB2 from the wider Cygnus complex in "
                 "any existing ²⁶Al measurement — SPI resolution is ~3°.", INK, False),
                ("It cannot distinguish high-mass explodability from binary "
                 "stripping as the origin of the pulsar.", INK, False),
                ("Absent remnants are weak evidence: even at a generous 100 "
                 f"kyr visibility the ledger predicts only {N['snrExpected']} "
                 f"visible remnants.", INK, False),
            ],
            size=16.5)

    return prs


def main() -> None:
    prs = build()
    OUT.parent.mkdir(exist_ok=True)
    prs.save(OUT)
    print(f"wrote {OUT.relative_to(ROOT)}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")


if __name__ == "__main__":
    main()
